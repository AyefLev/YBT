from dataclasses import dataclass
from pathlib import Path
import unicodedata

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


@dataclass(frozen=True)
class ParsedText:
    content: str
    page_no: int | None = None
    slide_no: int | None = None


SUPPORTED_SUFFIXES = {".txt", ".md", ".docx", ".pdf", ".pptx"}


class MaterialNeedsVisionError(RuntimeError):
    """Raised when a PDF has no trustworthy extractable text."""


def parse_material(path: Path, suffix: str) -> list[ParsedText]:
    if suffix in {".txt", ".md"}:
        return _parse_text(path)
    if suffix == ".docx":
        return _parse_docx(path)
    if suffix == ".pdf":
        return _parse_pdf(path)
    if suffix == ".pptx":
        return _parse_pptx(path)
    return []


def _parse_text(path: Path) -> list[ParsedText]:
    try:
        text = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        text = path.read_text(encoding="utf-8", errors="ignore")
    return [ParsedText(text)] if text.strip() else []


def _parse_docx(path: Path) -> list[ParsedText]:
    document = Document(path)
    parts: list[str] = []
    parts.extend(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" ".join(cells))
    text = "\n".join(parts)
    return [ParsedText(text)] if text.strip() else []


def _parse_pdf(path: Path) -> list[ParsedText]:
    reader = PdfReader(path)
    parsed: list[ParsedText] = []
    low_quality_pages = 0
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if not text.strip():
            low_quality_pages += 1
            continue
        if _is_low_quality_extracted_text(text):
            low_quality_pages += 1
            continue
        if text.strip():
            parsed.append(ParsedText(text, page_no=index))
    if len(reader.pages) > 0 and not parsed and low_quality_pages:
        raise MaterialNeedsVisionError("PDF 文本提取为空或质量过低，需要视觉模型解析。")
    return parsed


def _parse_pptx(path: Path) -> list[ParsedText]:
    presentation = Presentation(path)
    parsed: list[ParsedText] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        text = "\n".join(parts)
        if text.strip():
            parsed.append(ParsedText(text, slide_no=index))
    return parsed


def _is_low_quality_extracted_text(text: str) -> bool:
    normalized = text.strip()
    if not normalized:
        return True
    suspicious_chars = sum(1 for char in normalized if char in {"�", "\ufffd", "□", "\x00"})
    if suspicious_chars / max(len(normalized), 1) > 0.08:
        return True
    meaningful_chars = [char for char in normalized if not char.isspace()]
    control_chars = [
        char
        for char in meaningful_chars
        if unicodedata.category(char) in {"Cc", "Cf", "Cs", "Co", "Cn"}
    ]
    if len(control_chars) / max(len(meaningful_chars), 1) > 0.08:
        return True
    latin_extended_chars = [
        char
        for char in meaningful_chars
        if 0x0100 <= ord(char) <= 0x024F
    ]
    cjk_chars = [
        char
        for char in meaningful_chars
        if "\u4e00" <= char <= "\u9fff"
    ]
    if latin_extended_chars and len(cjk_chars) < 3 and len(latin_extended_chars) / max(len(meaningful_chars), 1) > 0.15:
        return True
    visible_chars = sum(1 for char in normalized if not char.isspace())
    if visible_chars < 8:
        return True
    return False
