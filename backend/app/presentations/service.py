from __future__ import annotations

import json
import re
from datetime import datetime, timezone
from pathlib import Path
from uuid import uuid4

from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sqlalchemy.orm import Session

from app.ai.prompts import POSTGRADUATE_AUDIENCE_REQUIREMENTS
from app.ai.schemas import AIResult
from app.ai.service import generate_text
from app.auth.models import User
from app.core.config import get_settings
from app.exports.models import ExportRecord
from app.lessons.models import Lesson
from app.presentations.schemas import (
    PresentationGenerateRequest,
    PresentationGenerateResponse,
    PresentationSlideRead,
)


def generate_lesson_presentation_file(
    db: Session,
    *,
    lesson: Lesson,
    payload: PresentationGenerateRequest,
    current_user: User,
) -> PresentationGenerateResponse:
    target_count = payload.slide_count or 6
    prompt = _build_presentation_prompt(lesson, payload=payload, target_count=target_count)
    provider_status: AIResult | None = generate_text(
        db,
        "presentation",
        prompt,
        user_id=current_user.id,
    )
    slides = _parse_slides(provider_status.content, target_count=target_count)
    if not slides:
        slides = _fallback_slides(lesson, target_count=target_count)

    file_path = _write_pptx(lesson=lesson, slides=slides)
    try:
        record = ExportRecord(
            resource_type="presentation",
            resource_id=lesson.id,
            file_path=str(file_path),
            exported_by=current_user.id,
        )
        db.add(record)
        db.flush()
    except Exception:
        _remove_file(file_path)
        raise

    return PresentationGenerateResponse(
        lesson_id=lesson.id,
        status="generated",
        queued=False,
        message="PPT generated.",
        slides=slides,
        provider_status=provider_status,
        export_id=record.id,
        download_url=f"/api/presentations/exports/{record.id}/pptx",
        filename=file_path.name,
    )


def _build_presentation_prompt(
    lesson: Lesson,
    *,
    payload: PresentationGenerateRequest,
    target_count: int,
) -> str:
    return "\n".join(
        [
            "Create a teaching PowerPoint outline as strict JSON.",
            f"slide_count={target_count}",
            f"title={lesson.title}",
            f"subject={lesson.subject}",
            f"chapter={lesson.chapter}",
            f"duration_minutes={lesson.duration_minutes}",
            f"style={payload.style or 'default'}",
            f"include_exercises={payload.include_exercises}",
            f"description={payload.description or 'Let the model choose the best slide plan.'}",
            "Return only JSON in this shape:",
            '{"slides":[{"title":"...","bullets":["..."],"speaker_notes":"...","visual_prompt":"..."}]}',
            "Keep each slide focused and suitable for classroom teaching.",
            "Design each slide as a class-ready postgraduate entrance exam lesson: concise title, progressive bullets, worked-example flow, and speaker notes.",
            "For visual_prompt, describe useful diagrams, formula layouts, tables, or board-writing cues instead of decorative images.",
            POSTGRADUATE_AUDIENCE_REQUIREMENTS,
            "Lesson content:",
            lesson.current_content[:8000],
        ]
    )


def _parse_slides(raw: str, *, target_count: int) -> list[PresentationSlideRead]:
    payload = _load_json_payload(raw)
    if payload is None:
        return []
    raw_slides = payload.get("slides") if isinstance(payload, dict) else payload
    if not isinstance(raw_slides, list):
        return []

    slides: list[PresentationSlideRead] = []
    for index, item in enumerate(raw_slides[:target_count], start=1):
        if not isinstance(item, dict):
            continue
        title = str(item.get("title") or f"Slide {index}").strip()
        bullets_value = item.get("bullets") or []
        if isinstance(bullets_value, str):
            bullets = [line.strip(" -*\t") for line in bullets_value.splitlines() if line.strip()]
        elif isinstance(bullets_value, list):
            bullets = [str(value).strip() for value in bullets_value if str(value).strip()]
        else:
            bullets = []
        slides.append(
            PresentationSlideRead(
                slide_no=len(slides) + 1,
                title=title[:120],
                bullets=bullets[:8],
                speaker_notes=str(item.get("speaker_notes") or "").strip()[:1200],
                visual_prompt=str(item.get("visual_prompt") or "").strip()[:600],
            )
        )
    return slides


def _load_json_payload(raw: str) -> object | None:
    raw = raw.strip()
    if not raw:
        return None
    candidates = [raw]
    object_start = raw.find("{")
    object_end = raw.rfind("}")
    if object_start >= 0 and object_end > object_start:
        candidates.append(raw[object_start : object_end + 1])
    array_start = raw.find("[")
    array_end = raw.rfind("]")
    if array_start >= 0 and array_end > array_start:
        candidates.append(raw[array_start : array_end + 1])

    for candidate in candidates:
        try:
            return json.loads(candidate)
        except json.JSONDecodeError:
            continue
    return None


def _fallback_slides(lesson: Lesson, *, target_count: int) -> list[PresentationSlideRead]:
    lines = [line.strip(" #*-") for line in lesson.current_content.splitlines() if line.strip()]
    if not lines:
        lines = [lesson.current_content.strip() or lesson.title]

    slides: list[PresentationSlideRead] = [
        PresentationSlideRead(
            slide_no=1,
            title=lesson.title,
            bullets=[
                f"Subject: {lesson.subject}",
                f"Chapter: {lesson.chapter}",
                f"Duration: {lesson.duration_minutes} minutes",
            ],
            speaker_notes="Open with lesson goals and expected outcomes.",
            visual_prompt="Clean title slide for a classroom lesson.",
        )
    ]
    remaining = max(0, target_count - 1)
    chunk_size = max(1, len(lines) // max(1, remaining))
    for index in range(remaining):
        chunk = lines[index * chunk_size : (index + 1) * chunk_size] or lines[:3]
        slides.append(
            PresentationSlideRead(
                slide_no=len(slides) + 1,
                title=chunk[0][:80] if chunk else f"Lesson Point {index + 1}",
                bullets=chunk[:5],
                speaker_notes="Explain this part with examples and quick checks.",
                visual_prompt="Simple teaching diagram related to the slide bullets.",
            )
        )
    return slides[:target_count]


def _write_pptx(*, lesson: Lesson, slides: list[PresentationSlideRead]) -> Path:
    export_dir = get_settings().export_dir
    export_dir.mkdir(parents=True, exist_ok=True)
    file_path = _presentation_file_path(export_dir, lesson)

    deck = PptxPresentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    for slide_data in slides:
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        _paint_slide(slide, slide_data=slide_data)
        if slide_data.speaker_notes:
            slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes

    deck.save(file_path)
    return file_path


def _paint_slide(slide, *, slide_data: PresentationSlideRead) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(248, 250, 252)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(37, 99, 235)
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(8.7), Inches(0.72))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.word_wrap = True
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = slide_data.title
    title_paragraph.font.size = Pt(30)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(15, 23, 42)

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), Inches(0.42), Inches(2.35), Inches(0.46))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(219, 234, 254)
    badge.line.fill.background()
    badge_frame = badge.text_frame
    badge_frame.clear()
    badge_frame.margin_left = Inches(0.08)
    badge_frame.margin_right = Inches(0.08)
    badge_paragraph = badge_frame.paragraphs[0]
    badge_paragraph.text = "考研课堂讲义"
    badge_paragraph.alignment = PP_ALIGN.CENTER
    badge_paragraph.font.size = Pt(13)
    badge_paragraph.font.bold = True
    badge_paragraph.font.color.rgb = RGBColor(30, 64, 175)

    content_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.55),
        Inches(1.35),
        Inches(7.7),
        Inches(5.35),
    )
    content_panel.fill.solid()
    content_panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_panel.line.color.rgb = RGBColor(226, 232, 240)
    content_frame = content_panel.text_frame
    content_frame.clear()
    content_frame.margin_left = Inches(0.32)
    content_frame.margin_right = Inches(0.25)
    content_frame.margin_top = Inches(0.22)
    content_frame.word_wrap = True
    heading = content_frame.paragraphs[0]
    heading.text = "课堂讲解"
    heading.font.size = Pt(16)
    heading.font.bold = True
    heading.font.color.rgb = RGBColor(37, 99, 235)
    bullets = slide_data.bullets or ["围绕本页主题展开核心概念、方法步骤和课堂提问。"]
    for bullet in bullets[:7]:
        paragraph = content_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_before = Pt(8)
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(30, 41, 59)

    visual_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.55),
        Inches(1.35),
        Inches(4.2),
        Inches(3.18),
    )
    visual_panel.fill.solid()
    visual_panel.fill.fore_color.rgb = RGBColor(239, 246, 255)
    visual_panel.line.color.rgb = RGBColor(147, 197, 253)
    visual_frame = visual_panel.text_frame
    visual_frame.clear()
    visual_frame.margin_left = Inches(0.24)
    visual_frame.margin_right = Inches(0.22)
    visual_frame.margin_top = Inches(0.2)
    visual_frame.word_wrap = True
    visual_title = visual_frame.paragraphs[0]
    visual_title.text = "视觉设计"
    visual_title.font.size = Pt(15)
    visual_title.font.bold = True
    visual_title.font.color.rgb = RGBColor(29, 78, 216)
    visual_body = visual_frame.add_paragraph()
    visual_body.text = slide_data.visual_prompt or "用板书、公式框或坐标示意图辅助解释本页内容。"
    visual_body.space_before = Pt(8)
    visual_body.font.size = Pt(16)
    visual_body.font.color.rgb = RGBColor(30, 64, 175)

    notes_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.55),
        Inches(4.75),
        Inches(4.2),
        Inches(1.95),
    )
    notes_panel.fill.solid()
    notes_panel.fill.fore_color.rgb = RGBColor(240, 253, 244)
    notes_panel.line.color.rgb = RGBColor(134, 239, 172)
    notes_frame = notes_panel.text_frame
    notes_frame.clear()
    notes_frame.margin_left = Inches(0.24)
    notes_frame.margin_right = Inches(0.22)
    notes_frame.margin_top = Inches(0.18)
    notes_frame.word_wrap = True
    notes_title = notes_frame.paragraphs[0]
    notes_title.text = "讲稿备注"
    notes_title.font.size = Pt(15)
    notes_title.font.bold = True
    notes_title.font.color.rgb = RGBColor(22, 101, 52)
    notes_body = notes_frame.add_paragraph()
    notes_body.text = slide_data.speaker_notes or "结合学生基础补充推导、例题和易错点。"
    notes_body.space_before = Pt(7)
    notes_body.font.size = Pt(14)
    notes_body.font.color.rgb = RGBColor(21, 128, 61)


def _presentation_file_path(export_dir: Path, lesson: Lesson) -> Path:
    safe_title = re.sub(r"[^A-Za-z0-9._-]+", "-", lesson.title).strip("-") or "lesson"
    for _ in range(10):
        timestamp = datetime.now(timezone.utc).strftime("%Y%m%d%H%M%S%f")
        candidate = export_dir / f"presentation-{lesson.id}-{safe_title}-{timestamp}-{uuid4().hex}.pptx"
        if not candidate.exists():
            return candidate
    raise RuntimeError("Could not allocate a unique presentation filename")


def _remove_file(file_path: Path) -> None:
    try:
        file_path.unlink(missing_ok=True)
    except OSError:
        pass
