import json
from io import BytesIO

from pptx import Presentation


def _auth_headers(client, username: str = "presentation_teacher") -> dict[str, str]:
    client.post(
        "/api/auth/register",
        json={
            "username": username,
            "email": f"{username}@example.com",
            "password": "secret-password",
            "display_name": "Presentation Teacher",
        },
    )
    login_response = client.post(
        "/api/auth/login",
        json={"username": username, "password": "secret-password"},
    )
    token = login_response.json()["access_token"]
    return {"Authorization": f"Bearer {token}"}


def _create_lesson(client, headers: dict[str, str]) -> int:
    response = client.post(
        "/api/lessons",
        headers=headers,
        json={
            "title": "Derivative Lesson",
            "subject": "Math",
            "chapter": "Derivatives",
            "stage": "Foundation",
            "duration_minutes": 45,
            "student_level": "Basic",
            "content": "Teaching goal: understand derivative definitions and solve simple tangent slope problems.",
        },
    )
    assert response.status_code == 201
    return int(response.json()["id"])


def test_generate_lesson_presentation_returns_slides_and_downloadable_pptx(client, monkeypatch):
    import httpx

    from app.cache.client import clear_cache_backend
    from app.core.config import get_settings

    captured: dict[str, object] = {}
    model_payload = {
        "slides": [
            {
                "title": "Derivative Definition",
                "bullets": ["Average rate of change", "Limit process", "Instantaneous slope"],
                "speaker_notes": "Connect the formula to a tangent line.",
                "visual_prompt": "Coordinate plane with secant becoming tangent.",
            },
            {
                "title": "Class Practice",
                "bullets": ["Compute f'(x)", "Interpret slope", "Discuss common mistakes"],
                "speaker_notes": "Leave two minutes for peer checking.",
                "visual_prompt": "Worked example with highlighted derivative steps.",
            },
        ]
    }

    def fake_post(url, *, headers, json, timeout):
        captured["url"] = url
        captured["authorization"] = headers["Authorization"]
        captured["model"] = json["model"]
        captured["prompt"] = json["messages"][0]["content"]
        return httpx.Response(
            200,
            request=httpx.Request("POST", url),
            json={
                "choices": [{"message": {"content": json_dumps(model_payload)}}],
                "usage": {"prompt_tokens": 11, "completion_tokens": 29},
            },
        )

    monkeypatch.setenv("LLM_API_KEY", "generate-key")
    monkeypatch.setenv("LLM_BASE_URL", "https://generate.example/v1")
    monkeypatch.setenv("LLM_MODEL", "generate-model")
    monkeypatch.setenv("LLM_MOCK_ON_FAILURE", "false")
    monkeypatch.setattr(httpx, "post", fake_post)
    get_settings.cache_clear()
    clear_cache_backend()
    headers = _auth_headers(client)
    lesson_id = _create_lesson(client, headers)

    response = client.post(
        f"/api/presentations/lesson/{lesson_id}/generate",
        headers=headers,
        json={
            "slide_count": 2,
            "description": "Focus on a concise visual explanation of derivatives.",
            "include_exercises": True,
        },
    )

    assert response.status_code == 200
    body = response.json()
    assert body["lesson_id"] == lesson_id
    assert body["status"] == "generated"
    assert body["queued"] is False
    assert body["provider_status"]["model"] == "generate-model"
    assert [slide["title"] for slide in body["slides"]] == [
        "Derivative Definition",
        "Class Practice",
    ]
    assert body["download_url"].startswith("/api/presentations/exports/")
    assert captured["url"] == "https://generate.example/v1/chat/completions"
    assert captured["authorization"] == "Bearer generate-key"
    assert "slide_count=2" in str(captured["prompt"])
    assert "Focus on a concise visual explanation of derivatives." in str(captured["prompt"])
    assert "postgraduate entrance exam lesson" in str(captured["prompt"])
    assert "面向考研备考人群" in str(captured["prompt"])

    download_response = client.get(body["download_url"], headers=headers)

    assert download_response.status_code == 200
    assert download_response.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert download_response.content.startswith(b"PK")

    deck = Presentation(BytesIO(download_response.content))
    slide_texts = [
        "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        for slide in deck.slides
    ]
    assert any("视觉设计" in text for text in slide_texts)
    assert any("课堂讲解" in text for text in slide_texts)


def json_dumps(value: object) -> str:
    return json.dumps(value, ensure_ascii=False)
